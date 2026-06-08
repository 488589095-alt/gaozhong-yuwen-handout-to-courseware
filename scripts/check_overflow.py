# -*- coding: utf-8 -*-
"""
check_overflow.py —— 课件出稿前自检（防出框三道防线的第②道，自动检测）

检测两类问题：
  1. 文字溢出：逐页逐文本框，按 字号×行距×1.08(字宽保守系数) 估算文字所需高度，
     超过框高即报嫌疑。系数 1.08 模拟真实中文字体（楷体/黑体/仿宋）比理论字宽略宽。
  2. shape 越界：内容元素（文本框/表格）超出 10×7.5in 页面。

为什么需要：渲染脚本的 _fit_size 已按实际行距缩字号（第①道防线），但估算总有误差；
本脚本是独立复核。第③道防线是 soffide 转 PDF 后全页缩略网格人工视觉扫。

用法：python3 check_overflow.py <课件.pptx>
退出码：0=通过；1=有嫌疑（必须修复后重跑，直到双 0 才能交付）
"""
import math
import sys

from pptx import Presentation


def check(path):
    p = Presentation(path)
    text_sus, bound_sus = [], []
    for i, s in enumerate(p.slides):
        for sh in s.shapes:
            # 模板装饰跳过（出血/画布外摆放是设计意图，如 T17 的 offscreen Shape 0/1）
            if sh.shape_type == 13 or sh.name.startswith(("Image", "图片", "组合", "Shape")):
                continue
            if not (sh.has_text_frame or sh.has_table):
                continue
            try:
                t, h = sh.top / 914400, sh.height / 914400
                l, w = sh.left / 914400, sh.width / 914400
            except TypeError:
                continue
            # ② shape 越界
            if t + h > 7.55 or t < -0.05 or l < -0.05 or l + w > 10.1:
                bound_sus.append((i + 1, sh.name, round(t + h, 2), round(l + w, 2)))
            # ① 文字溢出（严格估算）
            if not sh.has_text_frame:
                continue
            txt = sh.text_frame.text.strip()
            if not txt or len(txt) < 8:        # 短标签不会溢出
                continue
            usable_w = w - 0.15
            need = 0.0
            for para in sh.text_frame.paragraphs:
                ptxt = "".join(r.text for r in para.runs)
                if not ptxt:
                    need += 0.15
                    continue
                sz = None
                for r in para.runs:
                    if r.font.size:
                        sz = r.font.size.pt
                        break
                sz = sz or 18
                ls = para.line_spacing if isinstance(para.line_spacing, float) else 1.0
                cpl = max(1, int(usable_w * 72 / (sz * 1.08)))
                lines = max(1, math.ceil(len(ptxt) / cpl))
                need += lines * sz / 72.0 * max(ls, 1.0)
            if need > h + 0.05:
                text_sus.append((i + 1, sh.name, round(need, 2), round(h, 2), txt[:16]))

    print(f"文件: {path}  共 {len(p.slides)} 页")
    print(f"① 文字溢出嫌疑: {len(text_sus)}")
    for x in text_sus:
        print(f"   P{x[0]:>3} [{x[1]}] 需 {x[2]}in > 框 {x[3]}in | {x[4]!r}")
    print(f"② shape 越界: {len(bound_sus)}")
    for x in bound_sus:
        print(f"   P{x[0]:>3} [{x[1]}] bottom={x[2]} right={x[3]}")
    ok = not text_sus and not bound_sus
    print("✅ 双 0 通过" if ok else "❌ 有嫌疑，必须修复（缩字号/加框高/拆页）后重跑")
    return 0 if ok else 1


if __name__ == "__main__":
    if len(sys.argv) != 2:
        sys.exit("用法: python3 check_overflow.py <课件.pptx>")
    sys.exit(check(sys.argv[1]))
