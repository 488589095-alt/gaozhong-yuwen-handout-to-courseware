# -*- coding: utf-8 -*-
"""
build_gushi_pptx.py —— 高一语文「讲义 → 课件」生成脚本（针对 高一语文.pptx 模板）

设计：
  - 版式来源 = 高一语文.pptx（22 张样本页）。每个"角色"克隆指定样本页，再换文字/加元素。
  - 克隆带关系重映射：本模板每页有大量图片直接挂在 slide 上，朴素 deepcopy 会让
    r:embed 失效，必须把源页关系复制到新页并改写 rId。
  - 内容来源 = content.json（讲义派生）。换其他讲次只需替换 content.json。
  - 答案只取【答案】，不取【解析】；选择题答案页把正确选项整行标红。

★ 字体严格按标杆课件分布（实测）：
    诗=黑体 | 题干/出处/题型标签/各级标题=微软雅黑 | 答案=宋体(全红) |
    教材正文/表格body=楷体 | 序号=楷体 | 封面/目录/结束大字=模板主题字体(不覆盖)

用法：
  python3 build_gushi_pptx.py --content content.json --template "高一语文.pptx" \
      -o 古诗.pptx [--structure-only]
"""

import argparse
import copy
import json
import math
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
SHAPE_TAGS = ("sp", "pic", "grpSp", "graphicFrame", "cxnSp")
REF_ATTRS = [qn('r:embed'), qn('r:link'), qn('r:id'),
             qn('r:pict'), qn('r:dm'), qn('r:lo'), qn('r:qs'), qn('r:cs')]

# ── 颜色（对标杆）──
RED = "FF0000"          # 出处/题型标签/答案/选择正确项
DARKRED = "7D292D"      # 模块分隔标题
PINKRED = "EA4D54"      # 难度星级/内容总结标题（标杆亮红）
INK = "1A1A1A"          # 正文黑

# ── 字体（严格对标杆实测分布）──
YH = "微软雅黑"          # 出处/题型标签/各级标题
HEI = "黑体"            # 题目(题干)
SONG = "宋体"           # 答案
KAI = "楷体"            # 古诗正文 / 教材正文 / 表格body / 序号
FS = "仿宋"             # 古诗注释【注】

# ── 模板页索引（1-based，高一语文.pptx）──
T = dict(cover=1, toc=2, divider=3, table=8, legend=14,
         body=16, frame_ex=11, frame_pr=17, summary=16, end=18)


# ════════════════════════════════════════════════════════════
# 克隆（带关系重映射）
# ════════════════════════════════════════════════════════════
def clone_slide(prs, src):
    new = prs.slides.add_slide(src.slide_layout)
    spTree = new.shapes._spTree
    for ch in list(spTree):
        if etree.QName(ch).localname in SHAPE_TAGS:
            spTree.remove(ch)
    for el in src.shapes._spTree:
        if etree.QName(el).localname in SHAPE_TAGS:
            spTree.append(copy.deepcopy(el))
    relmap = {}
    for rId, rel in src.part.rels.items():
        if rel.reltype.endswith('slideLayout'):
            continue
        if rel.is_external:
            new_r = new.part.rels.get_or_add_ext_rel(rel.reltype, rel._target)
        else:
            new_r = new.part.relate_to(rel.target_part, rel.reltype)
        relmap[rId] = new_r
    for el in spTree.iter():
        for a in REF_ATTRS:
            v = el.get(a)
            if v in relmap:
                el.set(a, relmap[v])
    return new


# ════════════════════════════════════════════════════════════
# 文本 / 形状辅助
# ════════════════════════════════════════════════════════════
def _style_run(run, size=22, bold=False, color=INK, ea=YH, latin=None):
    """ea=None 时不覆盖字体，继承模板/主题（用于封面/目录/结束大字）。"""
    if size:
        run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor.from_string(color)
    if ea is None:
        return
    rpr = run._r.get_or_add_rPr()
    rpr.set("lang", "zh-CN")
    ea_el = rpr.find(f"{{{A_NS}}}ea")
    if ea_el is None:
        ea_el = etree.SubElement(rpr, f"{{{A_NS}}}ea")
    ea_el.set("typeface", ea)
    lat = latin or ea
    lat_el = rpr.find(f"{{{A_NS}}}latin")
    if lat_el is None:
        lat_el = etree.SubElement(rpr, f"{{{A_NS}}}latin")
    lat_el.set("typeface", lat)


def _fill_tf(tf, text, size=22, bold=False, color=INK, ea=YH,
             latin=None, align=None, line_spacing=None):
    tf.word_wrap = True
    tf.clear()
    first = True
    for line in str(text).split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if align is not None:
            p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        _style_run(run, size, bold, color, ea, latin)


def add_textbox(slide, text, left, top, width, height, **kw):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    anchor = kw.pop("anchor", None)
    if anchor is not None:
        tb.text_frame.vertical_anchor = anchor
    _fill_tf(tb.text_frame, text, **kw)
    return tb


def set_shape_text(slide, name, text, **kw):
    for sh in slide.shapes:
        if sh.name == name and sh.has_text_frame:
            _fill_tf(sh.text_frame, text, **kw)
            return sh
    return None


def clear_shapes(slide, names=None, prefixes=None):
    names = set(names or [])
    prefixes = tuple(prefixes or ())
    for sh in list(slide.shapes):
        nm = sh.name
        if nm in names or (prefixes and nm.startswith(prefixes)):
            sh._element.getparent().remove(sh._element)


def add_table(slide, header, rows, left, top, width, height,
              col_widths=None, head_size=18, body_size=18,
              head_ea=YH, body_ea=KAI):
    """表头微软雅黑深红底白字，正文楷体（贴标杆）。"""
    data = [header] + rows
    nr, nc = len(data), len(header)
    gf = slide.shapes.add_table(nr, nc, Inches(left), Inches(top),
                                Inches(width), Inches(height))
    tbl = gf.table
    if col_widths:
        # 列宽与列数不一致时容错：多则截断、少则按剩余宽度均分（信息类考情4列 vs 古诗5列）
        cw = list(col_widths[:nc])
        if len(cw) < nc:
            cw += [(width - sum(cw)) / (nc - len(cw))] * (nc - len(cw))
        for c, w in enumerate(cw):
            tbl.columns[c].width = Inches(w)
    for r in range(nr):
        for c in range(nc):
            cell = tbl.cell(r, c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            is_head = (r == 0)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(
                DARKRED if is_head else ("FFFFFF" if r % 2 else "FFF7EC"))
            _fill_tf(cell.text_frame, data[r][c],
                     size=head_size if is_head else body_size,
                     bold=is_head, color="FFFFFF" if is_head else INK,
                     ea=head_ea if is_head else body_ea,
                     align=PP_ALIGN.CENTER if (is_head or c != 1) else PP_ALIGN.LEFT)
    return gf


def add_tag(slide, text, left=0.3, top=1.27):
    """题型/测次标签：微软雅黑加粗红，如【必做题】【课前测】。"""
    return add_textbox(slide, text, left, top, 2.6, 0.45,
                       size=22, bold=True, color=RED, ea=YH,
                       anchor=MSO_ANCHOR.MIDDLE)


def add_module_badge(slide, text, left, top):
    """模块标题气泡（贴标杆）：白底红边、黑色微软雅黑粗字。"""
    w = 0.55 + len(text) * 0.42
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left),
                                Inches(top), Inches(w), Inches(0.55))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
    bg.line.color.rgb = RGBColor.from_string(DARKRED)
    bg.line.width = Pt(1.25)
    bg.shadow.inherit = False
    _fill_tf(bg.text_frame, text, size=22, bold=True, color="000000",
             ea=YH, align=PP_ALIGN.CENTER)
    bg.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return bg


CIRCLED = "①②③④⑤⑥⑦⑧⑨⑩⑪⑫"


def add_poem_textbox(slide, poem, left, top, width, height, size=22):
    """诗歌：黑体、居中；诗句中注释序号①②③设右上角上标（【注】行整体正常）。"""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in str(poem).split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.4
        is_note = line.lstrip().startswith("【注】")
        font = FS if is_note else KAI   # 诗句=楷体，注释=仿宋
        buf = ""

        def _flush(par=p, fnt=font):
            nonlocal buf
            if buf:
                r = par.add_run(); r.text = buf
                _style_run(r, size, False, INK, fnt); buf = ""
        for ch in line:
            if ch in CIRCLED and not is_note:
                _flush()
                r = p.add_run(); r.text = ch
                _style_run(r, max(12, int(size * 0.65)), False, INK, KAI)
                r._r.get_or_add_rPr().set("baseline", "30000")
            else:
                buf += ch
        _flush()
    return tb


def add_choice_block(slide, qtext, left, top, width, height, letter=None, size=22):
    """选择题题干：可把字母填入（ ）；若给 letter，正确选项整行标红（贴标杆）。"""
    shown = _fill_choice_letter(qtext, letter) if letter else qtext
    marker = f"（  {letter}  ）" if letter else None
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in shown.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = 1.3
        ls = line.lstrip()
        is_correct = bool(letter) and re.match(rf'^{letter}\s*[.．、]', ls)
        if marker and marker in line:          # 括号内答案字母标红加粗
            i = line.index(marker)
            segs = [(line[:i] + "（  ", INK, False), (letter, RED, True),
                    ("  ）" + line[i + len(marker):], INK, False)]
            for seg, col, b in segs:
                if not seg:
                    continue
                r = p.add_run(); r.text = seg
                _style_run(r, size, b, col, HEI)
        else:
            run = p.add_run(); run.text = line          # 题目=黑体；正确选项整行红
            _style_run(run, size, False, RED if is_correct else INK, HEI)
    return tb


def _fill_choice_letter(qtext, letter):
    for pat in [r"（\s+）", r"（\s*）", r"\(\s+\)", r"\(\s*\)"]:
        m = re.search(pat, qtext)
        if m:
            return qtext[:m.start()] + f"（  {letter}  ）" + qtext[m.end():]
    return qtext


def _est_h(text, cpl=25, line_h=0.42, pad=0.4):
    lines = sum(max(1, math.ceil(len(s) / cpl)) for s in str(text).split("\n"))
    return lines * line_h + pad


def _text_h(text, width_in, size, ls=1.3):
    """估算 text 在给定宽度/字号下的渲染高度（英寸）。"""
    cpl = max(8, int(width_in * 72 / (size * 1.08)))
    lines = sum(max(1, math.ceil(len(s) / cpl)) for s in str(text).split("\n"))
    return lines * size / 72.0 * ls


def _fit_size(text, width_in, avail_h, sizes=(22, 20, 18, 16), ls=1.25):
    for sz in sizes:
        cpl = max(8, int(width_in * 72 / (sz * 1.1)))   # 保守估每行字数（防出框）
        lines = sum(max(1, math.ceil(len(s) / cpl)) for s in str(text).split("\n"))
        if lines * (sz / 72.0 * ls) + 0.18 <= avail_h:
            return sz
    return sizes[-1]


# ════════════════════════════════════════════════════════════
# 角色渲染器
# ════════════════════════════════════════════════════════════
def r_cover(prs, C):
    s = clone_slide(prs, prs.slides[T["cover"] - 1])
    # 封面大字继承模板主题字体(ea=None)；标题用微软雅黑(标杆Text1)
    set_shape_text(s, "Text 0", C["lecture_no"], size=33.6, ea=None)
    set_shape_text(s, "Text 1", C["title"], size=44, bold=True, ea=YH,
                   line_spacing=1.1)
    set_shape_text(s, "Text 4", C["grade"], size=17.4, ea=None)
    set_shape_text(s, "Text 5", C["teacher"], size=17.4, ea=None)
    add_textbox(s, C.get("system_name", ""), 8.55, 0.16, 1.2, 0.2,
                size=8.4, color="FFFFFF", ea=None, align=PP_ALIGN.RIGHT)
    add_textbox(s, C.get("term", ""), 9.0, 0.31, 0.72, 0.2,
                size=10.2, color="FFFFFF", ea=None, align=PP_ALIGN.RIGHT)
    return s


def r_toc(prs, C):
    s = clone_slide(prs, prs.slides[T["toc"] - 1])
    clear_shapes(s, names=["Text 0", "Text 2", "Text 3", "Text 5",
                           "Text 6", "Text 8", "Text 9", "Text 11"],
                 prefixes=["Shape "])
    mods = C["modules"]
    y0, dy = 1.55, 0.86
    for i, name in enumerate(mods):
        y = y0 + i * dy
        circ = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.98),
                                  Inches(y), Inches(0.42), Inches(0.42))
        circ.fill.solid()
        circ.fill.fore_color.rgb = RGBColor.from_string("E8453C")
        circ.line.fill.background()
        circ.shadow.inherit = False
        _fill_tf(circ.text_frame, str(i + 1), size=20, bold=True,
                 color="FFFFFF", ea=YH, align=PP_ALIGN.CENTER)
        circ.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        add_textbox(s, name, 4.55, y - 0.04, 4.7, 0.5,
                    size=24, bold=True, color=INK, ea=YH,
                    anchor=MSO_ANCHOR.MIDDLE)
    return s


def r_divider(prs, no_cn, name):
    s = clone_slide(prs, prs.slides[T["divider"] - 1])
    clear_shapes(s, names=["Image 7"])   # 模板"空间几何体"标题贴图
    add_textbox(s, f"模块{no_cn}\n{name}", 3.93, 1.73, 3.83, 1.92,
                size=36, bold=True, color=DARKRED, ea=YH,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    return s


def r_table_page(prs, title, header, rows, col_widths, body_ea=KAI,
                 body_size=18, top=2.25, height=3.9):
    s = clone_slide(prs, prs.slides[T["table"] - 1])
    set_shape_text(s, "Text 0", "")
    set_shape_text(s, "Text 1", "")
    add_textbox(s, title, 3.3, 0.56, 3.4, 0.7, size=28, bold=True,
                color=DARKRED, ea=YH, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_table(s, header, rows, 0.5, top, 9.0, height,
              col_widths=col_widths, body_ea=body_ea, body_size=body_size)
    return s


def r_legend(prs, C):
    s = clone_slide(prs, prs.slides[T["legend"] - 1])
    set_shape_text(s, "Text 0", "")
    add_textbox(s, "难度星级说明", 4.15, 1.0, 5.0, 0.7, size=32, bold=True,
                color=PINKRED, ea=YH, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, "\n".join(C["star_legend"]), 2.6, 2.4, 5.2, 3.0,
                size=22, color=INK, ea=HEI, line_spacing=1.5,
                anchor=MSO_ANCHOR.MIDDLE)
    return s


def r_body(prs, title, paragraphs):
    s = clone_slide(prs, prs.slides[T["body"] - 1])
    clear_shapes(s, names=["Image 4"])   # 模板数学示例图
    set_shape_text(s, "Text 2", "")
    add_textbox(s, title, 4.1, 0.95, 5.2, 0.7, size=32, bold=True,
                color=DARKRED, ea=YH, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.MIDDLE)
    text = "\n".join(paragraphs)
    bsize = _fit_size(text, 8.2, 4.9, sizes=(22, 20, 18, 16, 14), ls=1.5)
    add_textbox(s, text, 0.9, 2.0, 8.2, 4.9, size=bsize,
                color=INK, ea=KAI, line_spacing=1.5)   # 正文楷体，按实际行距防出框
    return s


def _frame_slide(prs, frame, index=None, module_title=""):
    """诗/题/答 共用底；顶部模块标题气泡；清模板示例贴图与右上角课程徽标。"""
    if frame == "example":
        s = clone_slide(prs, prs.slides[T["frame_ex"] - 1])
        clear_shapes(s, names=["Image 2"])
        set_shape_text(s, "Text 2", "")
        if module_title:
            add_module_badge(s, module_title, 0.45, 0.42)
    else:
        s = clone_slide(prs, prs.slides[T["frame_pr"] - 1])
        clear_shapes(s, names=["Image 2", "Image 3", "Image 6", "Image 7",
                               "Image 8", "组合 2", "组合 21"])
        set_shape_text(s, "Text 2", index or "", size=22, bold=True,
                       color=INK, ea=KAI, align=PP_ALIGN.CENTER)   # 序号楷书
        set_shape_text(s, "Text 3", "")
        if module_title:
            add_module_badge(s, module_title, 1.0, 0.34)
    return s


def r_poem(prs, frame, source, poem, index=None, module_title=""):
    s = _frame_slide(prs, frame, index, module_title)
    # 出处微软雅黑红：按长度自动缩字号，确保一行放下不换行出框
    ssize = min(22, int(9.0 * 72 / (max(len(source), 1) * 1.06)))
    add_textbox(s, source, 0.4, 1.45, 9.1, 0.55, size=ssize, bold=True,
                color=RED, ea=YH, anchor=MSO_ANCHOR.MIDDLE)
    # 诗按"实际行距1.4"估算字号，留到页底 6.9（避免长诗如谢灵运 11 行溢出）
    psize = _fit_size(poem, 9.0, 4.85, sizes=(22, 20, 18, 16), ls=1.4)
    add_poem_textbox(s, poem, 0.5, 2.05, 9.0, 4.85, size=psize)
    return s


def _paginate_material(text):
    """信息/现代文类长阅读材料按"块"分页（对齐标杆：每页材料框用满、~10 行/约270-300字）。
    两条关键规则：
      ① 问答体把"答…"段绑定到前面的"问…"段成不可拆的块，使同一问答对落在同一页；
      ② 按**渲染行数**（非字数）累积分页——问答体段落多/短行多，按字数会低估行数致溢出。"""
    # 按"渲染行数"分页（问答体段落多/短行多，按字数会低估行数→溢出，必须按行数）
    cpl, max_lines = 27, 10        # 每行约27字(22pt,9.16in)；框高约4.95in≈10余行，留余量

    def nlines(t):
        return sum(max(1, math.ceil(len(x) / cpl)) for x in t.split("\n"))

    segs = str(text).split("\n")
    blocks = []
    for s in segs:
        if blocks and s.lstrip()[:1] == "答":     # “答…”并入前面的“问…”，问答成对
            blocks[-1] += "\n" + s
        else:
            blocks.append(s)
    pages, cur, cl = [], "", 0
    for b in blocks:
        bl = nlines(b)
        if bl > max_lines:                         # 单块(超长答/长论述段)超页→先结页再块内按句切
            if cur:
                pages.append(cur); cur, cl = "", 0
            while nlines(b) > max_lines:
                approx = max_lines * cpl
                cut = (b.rfind("。", 0, approx) + 1) or approx
                pages.append(b[:cut]); b = b[cut:]
            cur, cl = b, nlines(b)
        elif cur and cl + bl > max_lines:
            pages.append(cur); cur, cl = b, bl
        else:
            cur = (cur + "\n" + b) if cur else b
            cl += bl
    if cur:
        pages.append(cur)
    return pages or [""]


def r_material(prs, frame, source, material, index=None, module_title="", page=None, size=22):
    """信息/现代文阅读材料页：**楷体、统一字号 22pt**（不随页缩放，对齐标杆）、左对齐。
    长材料已按统一字号容量分页(_paginate_material)，故每页同字号、不溢出；
    首页带出处标签，多页在出处行右侧标"材料 n/m"。内容区起点固定→各页容量一致。"""
    s = _frame_slide(prs, frame, index, module_title)
    if source:
        ssize = min(22, int(9.0 * 72 / (max(len(source), 1) * 1.06)))
        add_textbox(s, source, 0.4, 1.45, 9.1, 0.55, size=ssize, bold=True,
                    color=RED, ea=YH, anchor=MSO_ANCHOR.MIDDLE)
    if page and page[1] > 1:
        add_textbox(s, f"材料 {page[0]}/{page[1]}", 7.7, 1.5, 1.9, 0.4,
                    size=14, color=INK, ea=YH, align=PP_ALIGN.RIGHT)
    top = 2.0                       # 内容区统一起点（不随有无出处变化→各页字号容量一致）
    add_textbox(s, material, 0.42, top, 9.16, 6.95 - top, size=size,
                color=INK, ea=KAI, line_spacing=1.4)         # 阅读材料=楷体·统一字号·用满高度
    return s


def r_question(prs, frame, qtext, index=None, tag=None, with_answer=False,
               answer=None, answer_letter=None, qtype=None, module_title=""):
    s = _frame_slide(prs, frame, index, module_title)
    qtop = 1.95 if tag else 1.77
    if tag:
        add_tag(s, tag)
    is_choice = (qtype == "choice") or (answer_letter is not None)
    if is_choice:
        h = 6.85 - qtop
        size = _fit_size(qtext, 9.2, h, sizes=(22, 20, 18), ls=1.3)  # 与实际行距一致
        letter = answer_letter if (with_answer and answer_letter) else None
        # 题干黑体；答案页括号字母红 + 正确选项整行标红
        add_choice_block(s, qtext, 0.4, qtop, 9.2, h, letter=letter, size=size)
    elif not with_answer:                       # 简答·题目空白页（题目=黑体）
        size = _fit_size(qtext, 9.2, 6.7 - qtop, sizes=(22, 20, 18), ls=1.3)
        add_textbox(s, qtext, 0.4, qtop, 9.2, 6.7 - qtop,
                    size=size, color=INK, ea=HEI, line_spacing=1.3)
    else:                                        # 简答·答案页：题干(黑体) + 答案(宋体全红)
        # 题干高度按内容自适应（补写题题干含长文段，简答题题干短），答案取剩余空间，
        # 给答案保底 ≥0.9in——避免长题干把答案区挤成负高度（信息类补写题踩过）
        total_h = 6.95 - qtop
        qsize = _fit_size(qtext, 9.2, total_h - 0.9, sizes=(22, 20, 18, 16), ls=1.3)
        qh = min(_text_h(qtext, 9.2, qsize, 1.3) + 0.12, total_h - 0.9)
        add_textbox(s, qtext, 0.4, qtop, 9.2, qh,
                    size=qsize, color=INK, ea=HEI, line_spacing=1.3)
        atop = qtop + qh + 0.2
        avail = 6.95 - atop
        asize = _fit_size(answer, 9.16, avail, sizes=(22, 20, 18, 16))
        add_textbox(s, answer, 0.42, atop, 9.16, avail, size=asize,
                    color=RED, ea=SONG, line_spacing=1.25)         # 答案宋体全红
    return s


# ════════════════════════════════════════════════════════════
# content → 有序 slide 规格 → 渲染
# ════════════════════════════════════════════════════════════
def build_specs(C):
    specs = []

    def add(role, summary, **payload):
        specs.append({"role": role, "summary": summary, **payload})

    add("cover", f"封面 {C['lecture_no']} {C['title']}")
    add("toc", "目录 " + "/".join(C["modules"]))

    add("divider", "模块一 学习目标", no_cn="一", name="学习目标")
    add("objective_table", "学习目标表")
    add("legend", "难度星级说明")

    add("divider", "模块二 考情分析", no_cn="二", name="考情分析")
    add("exam_table", "考情分析表")

    add("divider", "模块三 教材链接", no_cn="三", name="教材链接")
    add("body", "教材链接正文")

    add("divider", "模块四 典型例题", no_cn="四", name="典型例题")
    ex = C["example"]
    _emit_passage(add, "example", ex, None, "典型例题")   # 诗→poem / 阅读材料→material
    for q in ex["questions"]:
        _emit_qa(add, "example", q, None, tag=q.get("tag"), mt="典型例题")

    add("divider", "模块五 针对练习", no_cn="五", name="针对练习")
    for p in C["practices"]:
        idx = p["index"]
        _emit_passage(add, "practice", p, idx, "针对练习")
        for qi, q in enumerate(p["questions"]):
            _emit_qa(add, "practice", q, idx,
                     tag=(p.get("tag") if qi == 0 else None), mt="针对练习")

    add("summary", "内容总结")
    add("end", "结束页")
    return specs


def _emit_passage(add, frame, block, index, mt):
    """阅读对象：古诗(poem，1页) 或 信息/现代文长材料(material，自动分页)。"""
    if block.get("poem"):
        add("poem", "诗歌", frame=frame, index=index, module_title=mt,
            source=block["source"], poem=block["poem"])
    elif block.get("material"):
        pages = _paginate_material(block["material"])
        n = len(pages)
        for pi, frag in enumerate(pages):
            add("material", f"材料{pi+1}/{n}", frame=frame, index=index,
                module_title=mt, material=frag, page=(pi + 1, n),
                source=(block["source"] if pi == 0 else None))


def _emit_qa(add, frame, q, index, tag=None, mt=""):
    is_choice = q.get("type") == "choice"
    common = dict(frame=frame, index=index, qtext=q["text"],
                  qtype=q.get("type"), module_title=mt)
    add("question", "题目(空白)", tag=tag, with_answer=False, **common)
    if is_choice:
        add("answer", "答案(选择)", tag=tag, with_answer=True,
            answer_letter=q.get("answer_letter"), **common)
    else:
        add("answer", "答案(简答)", tag=tag, with_answer=True,
            answer=q.get("answer"), **common)


def render(content_path, template_path, out_path, structure_only=False):
    C = json.loads(Path(content_path).read_text(encoding="utf-8"))
    specs = build_specs(C)

    structure = {"template": Path(template_path).name,
                 "total_slides": len(specs),
                 "slides": [{"slide_num": i + 1, "role": s["role"],
                             "summary": s["summary"]}
                            for i, s in enumerate(specs)]}
    Path(out_path).with_name("slide_structure.json").write_text(
        json.dumps(structure, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"✅ slide_structure.json（{len(specs)} 页）")
    if structure_only:
        return

    prs = Presentation(str(template_path))
    orig = len(prs.slides)

    for sp in specs:
        role = sp["role"]
        if role == "cover":
            r_cover(prs, C)
        elif role == "toc":
            r_toc(prs, C)
        elif role == "divider":
            r_divider(prs, sp["no_cn"], sp["name"])
        elif role == "objective_table":
            o = C["objectives"]
            r_table_page(prs, "学习目标", o["header"], o["rows"],
                         col_widths=[1.1, 6.3, 1.6], body_ea=KAI, body_size=18)
        elif role == "exam_table":
            e = C["exam"]
            r_table_page(prs, "考情分析", e["header"], e["rows"],
                         col_widths=[0.9, 1.2, 2.5, 3.0, 1.4],
                         body_ea=KAI, body_size=12, top=1.8, height=4.6)
        elif role == "legend":
            r_legend(prs, C)
        elif role == "body":
            tb = C["textbook_link"]
            r_body(prs, tb["title"], tb["paragraphs"])
        elif role == "poem":
            r_poem(prs, sp["frame"], sp["source"], sp["poem"], sp.get("index"),
                   module_title=sp.get("module_title", ""))
        elif role == "material":
            r_material(prs, sp["frame"], sp.get("source"), sp["material"],
                       sp.get("index"), sp.get("module_title", ""), sp.get("page"))
        elif role in ("question", "answer"):
            r_question(prs, sp["frame"], sp["qtext"], index=sp.get("index"),
                       tag=sp.get("tag"), with_answer=sp["with_answer"],
                       answer=sp.get("answer"), answer_letter=sp.get("answer_letter"),
                       qtype=sp.get("qtype"), module_title=sp.get("module_title", ""))
        elif role == "summary":
            st = C.get("knowledge", {}).get("steps") or \
                ["第一步  释字义，破表层", "第二步  描景象，构画面",
                 "第三步  点手法，析技巧", "第四步  揭情感，悟主旨"]
            s = clone_slide(prs, prs.slides[T["summary"] - 1])
            clear_shapes(s, names=["Image 4"])
            set_shape_text(s, "Text 2", "内容总结", size=33, bold=True,
                           color=PINKRED, ea=YH)
            add_textbox(s, "炼字题解题步骤", 0.9, 1.7, 8.2, 0.6, size=26,
                        bold=True, color=RED, ea=YH, align=PP_ALIGN.CENTER)
            add_textbox(s, "\n".join(st), 2.2, 2.55, 5.8, 3.5, size=26,
                        color=INK, ea=YH, line_spacing=1.7,
                        anchor=MSO_ANCHOR.MIDDLE)
        elif role == "end":
            s = clone_slide(prs, prs.slides[T["end"] - 1])
            set_shape_text(s, "Text 0", C["end"]["big"], size=55, ea=None,
                           align=PP_ALIGN.CENTER, line_spacing=1.1)
            set_shape_text(s, "Text 3", C["end"]["small"], size=25.8, ea=None)

    for i in range(orig, 0, -1):
        sldIdLst = prs.slides._sldIdLst
        el = list(sldIdLst)[i - 1]
        prs.part.drop_rel(el.get(qn("r:id")))
        sldIdLst.remove(el)

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"✅ 课件已生成：{out_path}（{len(prs.slides)} 页）")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--content", required=True)
    ap.add_argument("--template", required=True)
    ap.add_argument("-o", "--output", required=True)
    ap.add_argument("--structure-only", action="store_true")
    a = ap.parse_args()
    render(a.content, a.template, a.output, a.structure_only)


if __name__ == "__main__":
    main()
